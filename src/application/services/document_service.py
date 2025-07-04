"""
文档处理服务
封装多种文档格式的上传、处理、向量化等核心业务逻辑
支持PDF、Word、Excel、PowerPoint、Markdown、文本等格式
"""

import os
import tempfile
import shutil
from pathlib import Path
from typing import Optional, Tuple, List
from datetime import datetime
import io
import base64
from PIL import Image

# 导入各种文档加载器
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredExcelLoader,
    UnstructuredPowerPointLoader,
    UnstructuredMarkdownLoader
)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain.schema import Document

# Google Gemini Vision API
import google.generativeai as genai

# Office文档处理
from docx import Document as DocxDocument
import pandas as pd
from pptx import Presentation

# 使用新的基础设施服务
from src.infrastructure.config.config_migration_adapter import get_legacy_config
from src.infrastructure.utilities import get_utility_service
from src.infrastructure import get_logger
from src.shared.state.application_state import app_state, FileInfo


class DocumentService:
    """文档处理服务 - 支持多种文档格式"""

    def __init__(self, model_service=None, config_service=None, logger_service=None, utility_service=None):
        """初始化文档服务

        Args:
            model_service: 模型管理服务实例，用于依赖注入
            config_service: 配置服务实例
            logger_service: 日志服务实例
            utility_service: 工具服务实例
        """
        self.model_service = model_service

        # 获取服务实例 (支持依赖注入)
        if config_service:
            # 如果提供了ConfigurationService，使用ConfigMigrationAdapter适配
            from src.infrastructure.config.config_migration_adapter import ConfigMigrationAdapter
            self.config = ConfigMigrationAdapter(config_service)
        else:
            self.config = get_legacy_config()
        self.logger = logger_service or get_logger()
        self.utility_service = utility_service or get_utility_service()

        # 支持的文档格式
        self.supported_formats = [".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md"]

        # 初始化文本分割器（支持语义分块）
        self._init_text_splitters()

        # 初始化嵌入模型
        self.embeddings = GoogleGenerativeAIEmbeddings(
            model="models/embedding-001",
            google_api_key=self.config.GOOGLE_API_KEY
        )

        # 初始化Google Gemini Vision API
        genai.configure(api_key=self.config.GOOGLE_API_KEY)
        self.vision_model = genai.GenerativeModel('gemini-1.5-flash')

        self.logger.info("DocumentService 初始化完成", extra={
            "chunk_size": self.config.CHUNK_SIZE,
            "chunk_overlap": self.config.CHUNK_OVERLAP,
            "use_semantic_chunking": self.config.USE_SEMANTIC_CHUNKING,
            "supported_formats": self.supported_formats
        })

    def _init_text_splitters(self):
        """初始化文本分割器"""
        # 传统分块器（备用）
        self.traditional_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.config.CHUNK_SIZE,
            chunk_overlap=self.config.CHUNK_OVERLAP,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )

        # 语义分块器
        try:
            from .semantic_text_splitter import AdaptiveSemanticSplitter
            self.semantic_splitter = AdaptiveSemanticSplitter(config=self.config)
            self.logger.info("语义分块器初始化成功")
        except ImportError as e:
            self.logger.warning(f"语义分块器导入失败: {e}，将使用传统分块器")
            self.semantic_splitter = None
        except Exception as e:
            self.logger.error(f"语义分块器初始化失败: {e}，将使用传统分块器")
            self.semantic_splitter = None

        # 根据配置选择默认分块器
        use_semantic = self.config.USE_SEMANTIC_CHUNKING
        if use_semantic and self.semantic_splitter:
            self.text_splitter = self.semantic_splitter
            self.logger.info("默认使用语义分块器")
        else:
            self.text_splitter = self.traditional_splitter
            self.logger.info("默认使用传统分块器")

    def process_document(self, file) -> str:
        """
        处理文档文件并创建向量数据库
        支持PDF、Word、Excel、PowerPoint、Markdown、文本等格式

        Args:
            file: 上传的文件对象

        Returns:
            处理结果消息
        """
        self.logger.info("开始处理文档文件", extra={"file": str(file)})

        if file is None:
            self.logger.warning("未提供文件")
            return "❌ 请选择一个文档文件"

        try:
            # 获取文件路径和名称
            file_path, file_name = self._get_file_info(file)

            self.logger.info("获取文件信息", extra={
                "file_path": file_path,
                "file_name": file_name
            })

            # 验证文件
            if not self._validate_file(file_path):
                return f"❌ 文件验证失败: {file_path}"

            self.logger.info("文件验证通过，开始处理文档")

            # 加载并处理文档
            documents = self._load_document(file_path)
            if not documents:
                return "❌ 文档文件为空或无法读取"

            self.logger.info("成功加载文档", extra={"pages": len(documents)})

            # 分割文档
            texts = self._split_documents(documents)
            self.logger.info("文档分割完成", extra={"chunks": len(texts)})

            # 检查是否有历史数据残留（用于用户提示）
            has_disk_data_before = self._check_vector_store_disk_state()

            # 创建向量存储
            success = self._create_vector_store(texts)
            if not success:
                return "❌ 向量数据库创建失败"

            # 更新文件记录
            self._update_file_record(file_name, len(documents), len(texts))

            self.logger.info("文档处理完成", extra={
                "file_name": file_name,
                "sections": len(documents),
                "chunks": len(texts)
            })

            # 构建返回消息
            result_message = f"✅ 文档 '{file_name}' 处理完成！\\n📄 文档段落: {len(documents)}\\n📝 文档片段: {len(texts)}"

            # 如果清理了历史数据，添加提示信息
            if has_disk_data_before and app_state.vectorstore is not None:
                result_message += "\\n🧹 已自动清理历史向量数据，确保检索结果基于当前会话文档"

            return result_message

        except Exception as e:
            error_msg = f"❌ 处理文档时发生错误: {str(e)}"
            self.logger.error("文档处理失败", exception=e, extra={"file": str(file)})
            return error_msg

    # 向后兼容方法
    def process_pdf(self, file) -> str:
        """
        向后兼容的PDF处理方法
        内部调用新的process_document方法
        """
        return self.process_document(file)

    def process_document_and_update_status(self, file, selected_model: str) -> Tuple[str, str, str, str]:
        """
        处理文档并更新系统状态

        Args:
            file: 上传的文件对象
            selected_model: 选择的模型

        Returns:
            (upload_status, model_status, system_status, file_list)
        """
        try:
            # 更新当前模型
            app_state.current_model = selected_model

            # 处理文档
            upload_status = self.process_document(file)

            # 获取更新后的状态
            model_status = f"当前模型: {app_state.current_model} (就绪)"
            system_status = self._get_system_status()
            file_list = self._get_uploaded_files_display()

            return upload_status, model_status, system_status, file_list

        except Exception as e:
            error_msg = f"❌ 处理失败: {str(e)}"
            return error_msg, "模型状态获取失败", "系统状态获取失败", "文件列表获取失败"

    # 向后兼容方法
    def process_pdf_and_update_status(self, file, selected_model: str) -> Tuple[str, str, str, str]:
        """
        向后兼容的PDF处理和状态更新方法
        内部调用新的process_document_and_update_status方法
        """
        return self.process_document_and_update_status(file, selected_model)

    def _get_file_info(self, file) -> Tuple[str, str]:
        """获取文件路径和名称"""
        if hasattr(file, 'name'):
            file_path = file.name
            file_name = Path(file_path).name
        else:
            file_path = str(file)
            file_name = Path(file_path).name
        return file_path, file_name

    def _validate_file(self, file_path: str) -> bool:
        """验证文件是否有效"""
        # 检查文件是否存在
        if not os.path.exists(file_path):
            self.logger.error(f"文件不存在: {file_path}")
            return False

        # 验证文件类型
        if not self.utility_service.validate_file_type(file_path, self.config.ALLOWED_FILE_TYPES):
            self.logger.error(f"不支持的文件类型: {file_path}")
            return False

        # 验证文件大小
        if not self.utility_service.validate_file_size(file_path, self.config.MAX_FILE_SIZE_MB):
            self.logger.error(f"文件大小超过限制: {file_path}")
            return False

        # 计算文件哈希（用于去重检测）
        file_hash = self.utility_service.calculate_file_hash(file_path)
        if file_hash:
            self.logger.debug(f"文件哈希计算完成: {file_path} -> {file_hash[:8]}...")

        return True

    def _load_document(self, file_path: str):
        """根据文件类型选择合适的加载器"""
        file_extension = Path(file_path).suffix.lower()

        self.logger.info(f"正在加载文档: {file_path}，文件类型: {file_extension}")

        try:
            documents = []

            if file_extension == ".pdf":
                # 使用PyPDFLoader处理PDF
                loader = PyPDFLoader(file_path)
                documents = loader.load()
                self.logger.info(f"PDF文档加载成功，共 {len(documents)} 页")

            elif file_extension == ".txt":
                # 使用TextLoader处理文本文件
                loader = TextLoader(file_path, encoding='utf-8')
                documents = loader.load()
                self.logger.info(f"文本文档加载成功，共 {len(documents)} 个文档")

            elif file_extension == ".md":
                # 使用UnstructuredMarkdownLoader处理Markdown
                try:
                    loader = UnstructuredMarkdownLoader(file_path)
                    documents = loader.load()
                    self.logger.info(f"Markdown文档加载成功，共 {len(documents)} 个段落")
                except Exception as e:
                    self.logger.warning(f"UnstructuredMarkdownLoader失败: {e}，尝试TextLoader")
                    loader = TextLoader(file_path, encoding='utf-8')
                    documents = loader.load()

            elif file_extension == ".docx":
                # Word文档处理
                try:
                    loader = UnstructuredWordDocumentLoader(file_path)
                    documents = loader.load()
                    self.logger.info(f"Word文档（Unstructured）加载成功，共 {len(documents)} 个段落")

                    # 检查是否有有效内容
                    total_content = "".join([doc.page_content.strip() for doc in documents])
                    if not total_content:
                        self.logger.warning("Word文档文字内容为空，尝试提取图片中的文字")

                        # 提取并识别图片中的文字
                        images = self._extract_images_from_docx(file_path)
                        if images:
                            self.logger.info(f"找到 {len(images)} 张图片，开始文字识别...")

                            all_text_content = []
                            for i, image in enumerate(images):
                                self.logger.info(f"正在识别第 {i+1} 张图片...")
                                text_content = self._recognize_text_from_image(image)
                                if text_content.strip():
                                    all_text_content.append(f"图片 {i+1} 内容:\n{text_content}")

                            if all_text_content:
                                combined_content = "\n\n".join(all_text_content)
                                documents = [Document(
                                    page_content=combined_content,
                                    metadata={"source": file_path, "extraction_method": "gemini_vision"}
                                )]
                                self.logger.info(f"图片文字识别完成，总内容长度: {len(combined_content)} 字符")
                            else:
                                self.logger.warning("图片中未识别到有效文字")
                        else:
                            self.logger.warning("Word文档中未找到图片")

                except Exception as e:
                    self.logger.warning(f"UnstructuredWordDocumentLoader失败: {e}，尝试python-docx直接读取")
                    try:
                        doc = DocxDocument(file_path)
                        content = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])

                        if content.strip():
                            documents = [Document(page_content=content, metadata={"source": file_path})]
                            self.logger.info(f"Word文档（python-docx）加载成功，内容长度: {len(content)} 字符")
                        else:
                            # 如果没有文字，尝试提取图片
                            self.logger.info("文档无文字内容，尝试提取图片...")
                            images = self._extract_images_from_docx(file_path)
                            if images:
                                all_text_content = []
                                for i, image in enumerate(images):
                                    text_content = self._recognize_text_from_image(image)
                                    if text_content.strip():
                                        all_text_content.append(f"图片 {i+1} 内容:\n{text_content}")

                                if all_text_content:
                                    combined_content = "\n\n".join(all_text_content)
                                    documents = [Document(
                                        page_content=combined_content,
                                        metadata={"source": file_path, "extraction_method": "gemini_vision"}
                                    )]
                                    self.logger.info(f"图片文字识别完成，总内容长度: {len(combined_content)} 字符")

                    except Exception as e2:
                        self.logger.error(f"python-docx也失败了: {e2}")

            elif file_extension == ".xlsx":
                # Excel文档处理
                try:
                    loader = UnstructuredExcelLoader(file_path)
                    documents = loader.load()
                    self.logger.info(f"Excel文档（Unstructured）加载成功，共 {len(documents)} 个表格")
                except Exception as e:
                    self.logger.warning(f"UnstructuredExcelLoader失败: {e}，尝试pandas直接读取")
                    try:
                        # 读取所有工作表
                        excel_file = pd.ExcelFile(file_path)
                        all_content = []

                        for sheet_name in excel_file.sheet_names:
                            df = pd.read_excel(file_path, sheet_name=sheet_name)
                            sheet_content = f"工作表: {sheet_name}\n{df.to_string(index=False)}"
                            all_content.append(sheet_content)

                        if all_content:
                            combined_content = "\n\n".join(all_content)
                            documents = [Document(page_content=combined_content, metadata={"source": file_path})]
                            self.logger.info(f"Excel文档（pandas）加载成功，内容长度: {len(combined_content)} 字符")

                    except Exception as e2:
                        self.logger.error(f"pandas读取Excel也失败了: {e2}")

            elif file_extension == ".pptx":
                # PowerPoint文档处理
                try:
                    loader = UnstructuredPowerPointLoader(file_path)
                    documents = loader.load()
                    self.logger.info(f"PowerPoint文档（Unstructured）加载成功，共 {len(documents)} 个幻灯片")
                except Exception as e:
                    self.logger.warning(f"UnstructuredPowerPointLoader失败: {e}，尝试python-pptx直接读取")
                    try:
                        prs = Presentation(file_path)
                        all_content = []

                        for i, slide in enumerate(prs.slides):
                            slide_content = f"幻灯片 {i+1}:\n"
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    slide_content += f"{shape.text}\n"
                            all_content.append(slide_content)

                        if all_content:
                            combined_content = "\n\n".join(all_content)
                            documents = [Document(page_content=combined_content, metadata={"source": file_path})]
                            self.logger.info(f"PowerPoint文档（python-pptx）加载成功，内容长度: {len(combined_content)} 字符")

                    except Exception as e2:
                        self.logger.error(f"python-pptx读取PowerPoint也失败了: {e2}")

            else:
                self.logger.error(f"不支持的文件格式: {file_extension}")
                return []

            # 验证文档内容
            for i, doc in enumerate(documents):
                content_length = len(doc.page_content.strip())
                self.logger.info(f"文档 {i+1} 内容长度: {content_length} 字符，内容预览: {doc.page_content[:100]}...")

                if content_length == 0:
                    self.logger.warning(f"文档 {i+1} 内容为空")

            return documents

        except Exception as e:
            self.logger.error(f"文档加载失败: {str(e)}")
            return []

    def _split_documents(self, documents):
        """智能分割文档（支持语义分块和传统分块）"""
        self.logger.info("正在分割文档...")

        # 检查配置和分块器可用性
        use_semantic = self.config.USE_SEMANTIC_CHUNKING
        fallback_enabled = self.config.FALLBACK_TO_TRADITIONAL

        texts = []
        split_method = "unknown"

        try:
            if use_semantic and self.semantic_splitter:
                # 尝试语义分块
                self.logger.info("尝试使用语义分块器...")
                texts = self.semantic_splitter.split_documents(documents, use_semantic=True)
                split_method = "semantic"

                # 验证语义分块结果
                if not texts and fallback_enabled:
                    self.logger.warning("语义分块结果为空，回退到传统分块")
                    texts = self.traditional_splitter.split_documents(documents)
                    split_method = "traditional_fallback"

            else:
                # 直接使用传统分块器
                self.logger.info("使用传统分块器...")
                texts = self.traditional_splitter.split_documents(documents)
                split_method = "traditional"

        except Exception as e:
            self.logger.error(f"分块过程出错: {str(e)}")

            if fallback_enabled and split_method != "traditional":
                self.logger.info("错误回退到传统分块...")
                try:
                    texts = self.traditional_splitter.split_documents(documents)
                    split_method = "traditional_error_fallback"
                except Exception as fallback_error:
                    self.logger.error(f"传统分块也失败: {str(fallback_error)}")
                    texts = []

        # 如果分割后没有文档片段，检查原文档内容
        if not texts and documents:
            self.logger.warning("文档分割后为空，检查原文档内容长度")
            for i, doc in enumerate(documents):
                content_length = len(doc.page_content.strip())
                self.logger.info(f"原文档 {i+1} 内容长度: {content_length} 字符")

                # 如果原文档内容太短，直接使用原文档
                if content_length > 0:
                    self.logger.info("使用原文档作为文档片段（内容较短）")
                    texts = documents
                    split_method = "no_split"
                    break

        # 记录分块统计信息
        chunk_stats = self._analyze_chunk_stats(texts)

        self.logger.info(f"文档分割完成", extra={
            "chunks": len(texts),
            "split_method": split_method,
            "avg_chunk_size": chunk_stats.get("avg_size", 0),
            "min_chunk_size": chunk_stats.get("min_size", 0),
            "max_chunk_size": chunk_stats.get("max_size", 0)
        })

        return texts

    def _analyze_chunk_stats(self, chunks) -> dict:
        """分析分块统计信息"""
        if not chunks:
            return {"avg_size": 0, "min_size": 0, "max_size": 0}

        sizes = [len(chunk.page_content) for chunk in chunks]

        return {
            "avg_size": sum(sizes) // len(sizes) if sizes else 0,
            "min_size": min(sizes) if sizes else 0,
            "max_size": max(sizes) if sizes else 0,
            "total_chunks": len(chunks)
        }

    def _create_vector_store(self, texts) -> bool:
        """创建或更新向量存储（增量模式）"""
        try:
            self.logger.info("正在处理向量数据库...")

            # 验证文档片段
            if not texts:
                self.logger.error("没有有效的文档片段用于创建向量数据库")
                return False

            # 过滤空内容的文档
            valid_texts = [doc for doc in texts if doc.page_content.strip()]

            if not valid_texts:
                self.logger.error("所有文档片段都是空的")
                return False

            self.logger.info(f"有效文档片段数量: {len(valid_texts)}")

            # 使用配置中的ChromaDB路径
            persist_directory = self.config.CHROMA_DB_PATH

                        # 检查应用状态和磁盘状态
            existing_vectorstore = app_state.vectorstore
            has_disk_data = self._check_vector_store_disk_state()

            if existing_vectorstore is None:
                # 应用状态中没有向量存储
                if has_disk_data:
                    # 检测到磁盘上有历史数据，先清理以确保干净状态
                    self.logger.warning("检测到历史向量数据残留，正在清理以确保干净状态...")
                    if not self._clear_vector_store():
                        self.logger.error("清理历史数据失败")
                        return False

                # 创建全新的向量数据库
                self.logger.info("创建全新向量数据库...")
                vectorstore = Chroma.from_documents(
                    documents=valid_texts,
                    embedding=self.embeddings,
                    persist_directory=persist_directory
                )

                # 设置到应用状态中
                app_state.vectorstore = vectorstore
                # 重置QA链以便使用新的向量存储
                app_state.qa_chain = None

                self.logger.info("向量数据库创建成功", extra={
                    "documents_added": len(valid_texts),
                    "operation": "create_clean",
                    "disk_cleaned": has_disk_data
                })
            else:
                # 增量添加到现有向量存储
                self.logger.info("向已有向量数据库增量添加文档...")
                try:
                    # 使用 add_documents 方法增量添加
                    existing_vectorstore.add_documents(valid_texts)

                    # 确保向量存储仍在应用状态中
                    app_state.vectorstore = existing_vectorstore
                    # 重置QA链以便刷新检索器
                    app_state.qa_chain = None

                    self.logger.info("文档成功增量添加到向量数据库", extra={
                        "documents_added": len(valid_texts),
                        "operation": "add"
                    })
                except Exception as add_error:
                    self.logger.error(f"增量添加失败: {str(add_error)}，尝试重新创建向量存储")
                    # 如果增量添加失败，尝试重新创建整个向量存储
                    vectorstore = Chroma.from_documents(
                        documents=valid_texts,
                        embedding=self.embeddings,
                        persist_directory=persist_directory
                    )

                    app_state.vectorstore = vectorstore
                    app_state.qa_chain = None

                    self.logger.warning("向量数据库已重新创建（增量添加失败后的回退操作）", extra={
                        "documents_added": len(valid_texts),
                        "operation": "recreate_fallback"
                    })

            return True

        except Exception as e:
            self.logger.error(f"向量数据库处理失败: {str(e)}")
            return False

    def _clear_vector_store(self) -> bool:
        """清理向量存储（删除持久化数据）"""
        try:
            persist_directory = self.config.CHROMA_DB_PATH
            persist_path = Path(persist_directory)

            if persist_path.exists():
                self.logger.info(f"检测到现有向量存储目录: {persist_directory}")

                # 安全删除目录及所有内容
                shutil.rmtree(persist_path)
                self.logger.info("向量存储目录已清理", extra={
                    "operation": "clear",
                    "directory": str(persist_path)
                })
            else:
                self.logger.info("向量存储目录不存在，无需清理")

            # 重置应用状态中的向量存储
            app_state.vectorstore = None
            app_state.qa_chain = None

            self.logger.info("向量存储状态已重置")
            return True

        except Exception as e:
            self.logger.error(f"清理向量存储失败: {str(e)}")
            return False

    def _check_vector_store_disk_state(self) -> bool:
        """检查磁盘上是否存在向量存储数据"""
        try:
            persist_directory = self.config.CHROMA_DB_PATH
            persist_path = Path(persist_directory)

            if not persist_path.exists():
                return False

            # 检查是否包含ChromaDB相关文件
            has_chroma_db = any([
                (persist_path / "chroma.sqlite3").exists(),
                any(persist_path.glob("*.sqlite3")),
                any(persist_path.iterdir())  # 任何文件或目录
            ])

            if has_chroma_db:
                self.logger.info("检测到磁盘上存在向量存储数据", extra={
                    "directory": str(persist_path),
                    "files": [f.name for f in persist_path.iterdir()]
                })

            return has_chroma_db

        except Exception as e:
            self.logger.error(f"检查向量存储磁盘状态失败: {str(e)}")
            return False

    def _extract_images_from_docx(self, docx_path: str) -> List[Image.Image]:
        """从Word文档中提取图片"""
        images = []
        try:
            doc = DocxDocument(docx_path)

            # 遍历文档中的所有关系（包括图片）
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        # 读取图片数据
                        image_data = rel.target_part.blob
                        # 转换为PIL Image
                        image = Image.open(io.BytesIO(image_data))
                        images.append(image)
                        self.logger.info(f"提取到图片，尺寸: {image.size}")
                    except Exception as e:
                        self.logger.warning(f"提取图片时出错: {e}")

        except Exception as e:
            self.logger.error(f"从Word文档提取图片失败: {e}")

        return images

    def _recognize_text_from_image(self, image: Image.Image) -> str:
        """使用Google Gemini Vision API识别图片中的文字"""
        try:
            # 将PIL图片转换为base64格式
            buffer = io.BytesIO()
            # 转换为RGB模式（去除alpha通道）
            if image.mode in ('RGBA', 'LA', 'P'):
                image = image.convert('RGB')
            image.save(buffer, format='JPEG', quality=95)
            image_data = buffer.getvalue()

            # 调用Gemini Vision API
            prompt = """
            请提取这张图片中的所有文字内容。
            要求：
            1. 保持原有的格式和结构
            2. 如果有表格，请用合适的格式表示
            3. 如果有列表，请保持列表格式
            4. 只返回文字内容，不要添加说明
            """

            response = self.vision_model.generate_content([
                prompt,
                {
                    'mime_type': 'image/jpeg',
                    'data': image_data
                }
            ])

            text_content = response.text if response.text else ""
            self.logger.info(f"Gemini Vision 识别到文字长度: {len(text_content)} 字符")

            return text_content

        except Exception as e:
            self.logger.error(f"Gemini Vision API调用失败: {e}")
            return ""

    def _update_file_record(self, file_name: str, pages: int, chunks: int):
        """更新文件记录"""
        file_info = FileInfo(
            name=file_name,
            upload_time=datetime.now(),
            pages=pages,
            chunks=chunks,
            model=app_state.current_model
        )
        app_state.add_uploaded_file(file_info)

    def _get_system_status(self) -> str:
        """获取系统状态"""
        try:
            state_info = app_state.get_state_info()

            status_parts = [
                f"🤖 当前模型: {state_info['current_model']}",
                f"📊 状态: {'就绪' if state_info['qa_chain_initialized'] else '未加载'}",
                f"📚 已上传文档: {state_info['uploaded_files_count']} 个",
                f"🔧 向量数据库: {'已初始化' if state_info['vectorstore_initialized'] else '未初始化'}",
                f"🕐 最后更新: {datetime.fromisoformat(state_info['last_update']).strftime('%H:%M:%S')}"
            ]

            return "\\n".join(status_parts)

        except Exception as e:
            return f"❌ 获取系统状态失败: {str(e)}"

    def _get_uploaded_files_display(self) -> str:
        """获取上传文件列表显示"""
        try:
            files = app_state.get_uploaded_files()

            if not files:
                return "暂无已上传的文件"

            file_list = ["## 📁 已上传的文件\\n"]

            for i, file_info in enumerate(files, 1):
                upload_time = file_info.upload_time.strftime("%Y-%m-%d %H:%M:%S")
                file_list.append(
                    f"**{i}. {file_info.name}**\\n"
                    f"   - 📅 上传时间: {upload_time}\\n"
                    f"   - 📄 页数: {file_info.pages}\\n"
                    f"   - 📝 文档片段: {file_info.chunks}\\n"
                    f"   - 🤖 处理模型: {file_info.model}\\n"
                )

            return "\\n".join(file_list)

        except Exception as e:
            return f"❌ 获取文件列表失败: {str(e)}"

    def get_uploaded_files_count(self) -> int:
        """获取已上传文件数量"""
        return len(app_state.get_uploaded_files())

    def clear_uploaded_files(self):
        """清空上传文件记录"""
        app_state.clear_uploaded_files()

    def clear_all_documents_and_storage(self) -> str:
        """清空所有文档记录和向量存储"""
        try:
            # 清空上传文件记录
            self.clear_uploaded_files()

            # 清理向量存储
            if self._clear_vector_store():
                self.logger.info("所有文档和向量存储已清理")
                return "✅ 已清空所有文档记录和向量存储"
            else:
                return "❌ 清理向量存储时发生错误"

        except Exception as e:
            error_msg = f"❌ 清理操作失败: {str(e)}"
            self.logger.error("清理所有数据失败", exception=e)
            return error_msg